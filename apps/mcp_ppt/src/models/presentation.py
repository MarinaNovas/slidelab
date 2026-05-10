from dataclasses import dataclass
from pptx import Presentation
from typing import Any, Literal, Optional
from pydantic import BaseModel, Field, model_validator

@dataclass
class PresentationModel:
    id: str
    title: str
    prs: Presentation
    template_name: str | None = None


class PresentationMetadata(BaseModel):
    title: str
    subtitle: Optional[str] = None
    template_name: Optional[str] = None


class ImageSpec(BaseModel):
    prompt: str
    style: Optional[str] = "business presentation illustration"
    width_ratio: int = 1
    height_ratio: int = 1
    seed: Optional[int] = None


class BaseSlide(BaseModel):
    type: Literal["section", "content", "image_content", "comparison_table"]
    title: str


class SectionSlide(BaseSlide):
    type: Literal["section"] = "section"


class ContentSlide(BaseSlide):
    type: Literal["content"] = "content"
    content: list[str] = Field(default_factory=list)


class ImageContentSlide(BaseSlide):
    type: Literal["image_content"] = "image_content"
    section_title: Optional[str] = None
    subtitle: Optional[str] = None
    content: list[str] = Field(default_factory=list)
    image: ImageSpec


class ComparisonTableSlide(BaseSlide):
    type: Literal["comparison_table"] = "comparison_table"
    sidebar_items: list[str] = Field(default_factory=list)
    table_title: Optional[str] = ""
    headers: list[str]
    rows: list[list[str]]


class PresentationPlan(BaseModel):
    metadata: PresentationMetadata
    agenda: list[str] = Field(default_factory=list, max_length=6)
    slides: list[SectionSlide | ContentSlide | ImageContentSlide | ComparisonTableSlide]
    add_thank_you: bool = True
    save: bool = True

    @model_validator(mode="after")
    def validate_plan(self):
        if not self.slides:
            raise ValueError("slides must contain at least one slide")

        image_slides = [slide for slide in self.slides if slide.type == "image_content"]
        main_slides = [slide for slide in self.slides if slide.type != "section"]

        if main_slides:
            image_ratio = len(image_slides) / len(main_slides)
            if image_ratio < 0.3:
                # Не падаем жёстко, но можно сделать warning в результате
                pass

        return self