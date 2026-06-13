from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER


class TemplateInspector:
    @staticmethod
    def inspect(template_path: Path) -> dict:
        if not template_path.exists():
            return {
                "status": "error",
                "message": f"Template not found: {template_path}",
            }

        prs = Presentation(str(template_path))
        layouts = []

        for master_index, master in enumerate(prs.slide_masters):
            print(f"\nMASTER {master_index}")
            for layout_index, layout in enumerate(master.slide_layouts):
                print(f"{layout_index} {layout.name}")
                placeholders = TemplateInspector._inspect_placeholders(layout)
                semantic_type = TemplateInspector._guess_layout_semantic_type(
                    layout_name = layout.name,
                    placeholders = placeholders,
                )
                layouts.append({
                    "master_index": master_index,
                    "layout_index": layout_index,
                    "layout_name": layout.name,
                    "semantic_type": semantic_type,
                    "placeholders": placeholders,
                })

        return {
            "status": "ok",
            "template_path": str(template_path),
            "masters_count": len(prs.slide_masters),
            "layouts_count": len(layouts),
            "layouts": layouts,
            "semantic_profile": TemplateInspector._build_semantic_profile(layouts),
        }

    @staticmethod
    def _map_placeholders_by_role(placeholders: list[dict]) -> dict:
        result = {}

        titles = [
            p for p in placeholders
            if p["type_name"] in ["TITLE", "CENTER_TITLE"]
        ]
        subtitles = [p for p in placeholders if p["type_name"] == "SUBTITLE"]
        bodies = [p for p in placeholders if p["type_name"] == "BODY"]
        objects = [p for p in placeholders if p["type_name"] == "OBJECT"]
        pictures = [p for p in placeholders if p["type_name"] == "PICTURE"]

        if titles:
            result["title"] = titles[0]["idx"]

        if subtitles:
            result["subtitle"] = subtitles[0]["idx"]

        if bodies:
            result["body"] = bodies[0]["idx"]

        if objects:
            result["content"] = objects[0]["idx"]

        if pictures:
            result["image"] = pictures[0]["idx"]

        if len(objects) >= 2:
            sorted_objects = sorted(objects, key = lambda p: p["left"])
            result["left_content"] = sorted_objects[0]["idx"]
            result["right_content"] = sorted_objects[1]["idx"]

        if len(bodies) >= 2:
            sorted_bodies = sorted(bodies, key = lambda p: p["left"])
            result["left_body"] = sorted_bodies[0]["idx"]
            result["right_body"] = sorted_bodies[1]["idx"]

        return result

    @staticmethod
    def _build_semantic_profile(layouts: list[dict]) -> dict:
        profile: dict[str, list[dict]] = {}

        for layout in layouts:
            semantic_type = layout["semantic_type"]

            if semantic_type == "unknown":
                continue

            layout_profile = {
                "master_index": layout["master_index"],
                "layout_index": layout["layout_index"],
                "layout_name": layout["layout_name"],
                "placeholders": TemplateInspector._map_placeholders_by_role(
                    layout["placeholders"]
                ),
            }
            profile.setdefault(semantic_type, []).append(layout_profile)
        return profile

    @staticmethod
    def _placeholder_type_name(ph_type) -> str:
        mapping = {
            PP_PLACEHOLDER.TITLE: "TITLE",
            PP_PLACEHOLDER.CENTER_TITLE: "CENTER_TITLE",
            PP_PLACEHOLDER.SUBTITLE: "SUBTITLE",
            PP_PLACEHOLDER.BODY: "BODY",
            PP_PLACEHOLDER.OBJECT: "OBJECT",
            PP_PLACEHOLDER.PICTURE: "PICTURE",
            PP_PLACEHOLDER.TABLE: "TABLE",
            PP_PLACEHOLDER.CHART: "CHART",
            PP_PLACEHOLDER.FOOTER: "FOOTER",
            PP_PLACEHOLDER.SLIDE_NUMBER: "SLIDE_NUMBER",
            PP_PLACEHOLDER.DATE: "DATE",
        }

        return mapping.get(ph_type, str(ph_type))

    @staticmethod
    def _guess_layout_semantic_type(
            layout_name: str,
            placeholders: list[dict],
    ) -> str:
        name = layout_name.lower()
        types = [p["type_name"] for p in placeholders]

        title_count = types.count("TITLE") + types.count("CENTER_TITLE")
        body_count = types.count("BODY")
        object_count = types.count("OBJECT")
        picture_count = types.count("PICTURE")
        print(f"{name=}" )
        if "agenda" in name:
            return "agenda"

        if "thank" in name or "closing" in name or "final" in name:
            return "thank_you"

        if "comparison" in name:
            return "comparison"

        if "two content" in name or object_count >= 2:
            return "two_columns"

        if picture_count >= 1:
            return "image_content"

        if "content" in name:
            return "content"

        if "section" in name:
            return "section"

        if "cover" in name or "title" in name or "first" in name:
            return "cover"

        if title_count >= 1 and (body_count >= 1 or object_count >= 1):
            return "content"

        return "unknown"

    @staticmethod
    def _inspect_placeholders(layout) -> list[dict]:
        placeholders = []

        for shape in layout.placeholders:
            ph_format = shape.placeholder_format
            ph_type = ph_format.type
            print(f"{ph_format.type}")

            placeholders.append({
                "idx": ph_format.idx,
                "type": str(ph_type),
                "type_name": TemplateInspector._placeholder_type_name(ph_type),
                "name": shape.name,
                "left": shape.left,
                "top": shape.top,
                "width": shape.width,
                "height": shape.height,
            })
        return placeholders





