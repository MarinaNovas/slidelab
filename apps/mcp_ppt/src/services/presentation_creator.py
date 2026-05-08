import re
import uuid
from pathlib import Path

from pptx import Presentation

from src.config import DEFAULT_TEMPLATE, TEMPLATES_DIR
from src.models.presentation import PresentationModel
from src.services.presentation_store import PresentationStore

class PresentationCreator:
    def __init__(self, store: PresentationStore) -> None:
        self.store = store

    def create(self, title: str, template_name: str | None = None) -> str:
        prs_id = self._generate_id(title)
        template_path = self._resolve_template_path(template_name)
        if template_path.exists():
            prs = Presentation(str(template_path))
        else:
            prs = Presentation()
        presentation = PresentationModel(
            id = prs_id,
            title = title,
            prs = prs,
        )
        self.store.add(presentation)
        return prs_id

    def get_info(self, prs_id: str) -> str:
        presentation = self.store.get(prs_id)
        prs = presentation.prs

        info = f"Presentation: {presentation.id}\n"
        info += f"Title: {presentation.title}\n"
        info += f"Number of slides: {len(prs.slides)}\n"

        return info

    def _generate_id(self, title: str) -> str:
        cleaned_title = re.sub(r"[^\w\s\-]", "", title)
        cleaned_title = re.sub(r"[\s\-]+", "_", cleaned_title)
        cleaned_title = cleaned_title.strip("_").lower()
        cleaned_title = cleaned_title[:20].rstrip("_")

        short_uuid = uuid.uuid4().hex[:8]

        if cleaned_title:
            return f"{cleaned_title}_{short_uuid}"

        return f"presentation_{short_uuid}"

    def _resolve_template_path(self, template_name: str | None) -> Path:
        if template_name:
            return TEMPLATES_DIR / template_name

        return DEFAULT_TEMPLATE
