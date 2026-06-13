from datetime import date
from pathlib import Path

from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from src.models.slide import (
    AgendaSlideData, ComparisonTableSlideData, ImageSlideData, TitleSlideData,
    ContentSlideData,
    SectionSlideData,
    TableSlideData,
)
from src.services.presentation_store import PresentationStore


class SlideCreator:
    def __init__(self, store: PresentationStore) -> None:
        self.store = store

    @staticmethod
    def _add_semantic_slide(presentation, semantic_type: str):
        profile = presentation.template_profile
        if not profile or semantic_type not in profile:
            raise ValueError(f"Layout '{semantic_type}' not found in template profile")
        layout_info = profile[semantic_type][0]
        prs = presentation.prs
        master = prs.slide_masters[layout_info["master_index"]]
        layout = master.slide_layouts[layout_info["layout_index"]]
        slide = prs.slides.add_slide(layout)
        return slide, layout_info["placeholders"]

    def add_title_slide(self, prs_id: str, data: TitleSlideData) -> int:
        presentation = self.store.get(prs_id)
        prs = presentation.prs

        slide, ph = SlideCreator._add_semantic_slide(
            presentation = presentation,
            semantic_type = "cover",
        )
        subtitle_idx = ph.get("subtitle") or ph.get("body") or ph.get("left-body")
        if subtitle_idx is not None and data.subtitle:
            subtitle_ph = self._get_placeholder(slide, subtitle_idx)
            subtitle_ph.text = data.subtitle

        title_ph = self._get_placeholder(slide, ph["title"])
        title_ph.text = data.title

        # если в шаблоне есть второй BODY — используем его под дату
        date_idx = ph.get("second_body") or ph.get("date") or ph.get("right_body")

        if date_idx is not None:
            date_ph = self._get_placeholder(slide, date_idx)
            date_ph.text = str(date.today())

        return len(prs.slides)

    def add_content_slide(self, prs_id: str, data: ContentSlideData) -> int:
        presentation = self.store.get(prs_id)
        prs = presentation.prs

        slide, ph = SlideCreator._add_semantic_slide(
            presentation = presentation,
            semantic_type = "content",
        )

        title_ph = self._get_placeholder(slide, ph["title"])

        content_idx = ph.get("content") or ph.get("body")
        body_ph = self._get_placeholder(slide, content_idx)

        title_ph.text = data.title
        body_ph.text = "\n".join(data.content)

        return len(prs.slides)

    def add_image_content_slide(self, prs_id: str, data: ImageSlideData) -> int:
        presentation = self.store.get(prs_id)
        prs = presentation.prs

        slide, ph = SlideCreator._add_semantic_slide(
            presentation = presentation,
            semantic_type = "image_content"
        )

        title_ph = self._get_placeholder(slide, ph["title"])
        body_ph = self._get_placeholder(slide, ph["body"])
        object_ph = self._get_placeholder(slide, ph["content"])
        picture_ph = self._get_placeholder(slide, ph["image"])

        title_ph.text = data.title
        body_ph.text = data.subtitle or ""

        print(f"{data.content}=")
        tf = object_ph.text_frame
        tf.clear()  # очищаем дефолтный текст

        for i, item in enumerate(data.content):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            p.text = item
            p.level = 0  # уровень bullet (0 = основной)

        if data.image_path:
            image_path = Path(data.image_path) #2bf37c60d1ae44c6b9f495035cc40eb0

            if not image_path.exists():
                raise ValueError(f"Image not found: {image_path}")

            picture_ph.insert_picture(str(image_path))

        return len(prs.slides)

    def add_section_slide(self, prs_id: str, data: SectionSlideData) -> int:
        presentation = self.store.get(prs_id)
        prs = presentation.prs

        slide, ph = SlideCreator._add_semantic_slide(
            presentation = presentation,
            semantic_type = "section",
        )
        title_ph = self._get_placeholder(slide, ph["title"])
        title_ph.text = data.section_title

        return len(prs.slides)

    def add_table_slide(self, prs_id: str, data: TableSlideData) -> int:
        presentation = self.store.get(prs_id)
        prs = presentation.prs

        slide = prs.slides.add_slide(prs.slide_layouts[13])

        title_shape = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(9), Inches(1)
        )
        title_shape.text = data.title
        title_shape.text_frame.paragraphs[0].font.size = Pt(28)
        title_shape.text_frame.paragraphs[0].font.bold = True

        rows_count = len(data.rows) + 1
        cols_count = len(data.headers)

        table = slide.shapes.add_table(
            rows_count,
            cols_count,
            Inches(1),
            Inches(2),
            Inches(8),
            Inches(0.5 * rows_count),
        ).table

        for index, header in enumerate(data.headers):
            cell = table.cell(0, index)
            cell.text = header
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        for row_index, row in enumerate(data.rows, start=1):
            for col_index, value in enumerate(row):
                if col_index < cols_count:
                    table.cell(row_index, col_index).text = str(value)

        return len(prs.slides)

    def add_thank_you_slide(self, prs_id: str) -> int:
        presentation = self.store.get(prs_id)

        if presentation is None:
            raise ValueError(f"Presentation '{prs_id}' not found")

        prs = presentation.prs
        slide, _ = SlideCreator._add_semantic_slide(
            presentation = presentation,
            semantic_type = "thank_you",
        )
        return len(prs.slides)

    def add_agenda_slide(self, prs_id: str, data: AgendaSlideData) -> int:
        presentation = self.store.get(prs_id)

        if presentation is None:
            raise ValueError(f"Presentation '{prs_id}' not found")

        prs = presentation.prs

        slide, ph = SlideCreator._add_semantic_slide(
            presentation = presentation,
            semantic_type = "agenda",
        )

        title_idx = ph.get("title") or ph.get("body")

        if title_idx is None:
            raise ValueError("Title placeholder not found for semantic layout 'agenda'")

        title_ph = self._get_placeholder(slide, title_idx)
        title_ph.text = data.title or "AGENDA"

        number_placeholder_ids = [44, 45, 46, 34, 35, 36]
        title_placeholder_ids = [59, 60, 56, 54, 55, 61]

        for i, item in enumerate(data.items[:6]):
            number_ph = self._get_placeholder(slide, number_placeholder_ids[i])
            item_title_ph = self._get_placeholder(slide, title_placeholder_ids[i])

            number_ph.text = str(item.number or "")
            item_title_ph.text = str(item.title or "")

        return len(prs.slides)

    def add_comparison_table_slide(
            self,
            prs_id: str,
            data: ComparisonTableSlideData,
    ) -> int:
        presentation = self.store.get(prs_id)

        if presentation is None:
            raise ValueError(f"Presentation '{prs_id}' not found")

        prs = presentation.prs
        slide, ph = SlideCreator._add_semantic_slide(
            presentation = presentation,
            semantic_type = "comparison",
        )

        title_ph = self._get_placeholder(slide, ph["title"])
        title_ph.text = data.title or ""

        table_title_idx = ph.get("body")
        sidebar_idx = ph.get("left_body")
        object_idx = ph.get("content")

        title_ph.text = data.title or ""

        if table_title_idx is not None:
            table_title_ph = self._get_placeholder(slide, table_title_idx)
            table_title_ph.text = data.table_title or ""

        if sidebar_idx is not None:
            sidebar_ph = self._get_placeholder(slide, sidebar_idx)
            sidebar_ph.text = "\n".join(
                str(item)
                for item in data.sidebar_items
                if item is not None
            )

        if object_idx is None:
            raise ValueError("Table placeholder not found for semantic layout 'comparison'")

        object_ph = self._get_placeholder(slide, object_idx)

        rows_count = len(data.rows) + 1
        cols_count = len(data.headers)

        graphic_frame = slide.shapes.add_table(
            rows_count,
            cols_count,
            object_ph.left,
            object_ph.top,
            object_ph.width,
            object_ph.height,
        )

        table = graphic_frame.table

        # header row
        for col_idx, header in enumerate(data.headers):
            cell = table.cell(0, col_idx)
            cell.text = str(header)

            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                for run in paragraph.runs:
                    run.font.bold = True
                    run.font.size = Pt(12)

        # body rows
        for row_idx, row in enumerate(data.rows, start = 1):
            for col_idx, value in enumerate(row):
                cell = table.cell(row_idx, col_idx)
                cell.text = str(value)

                for paragraph in cell.text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.LEFT
                    for run in paragraph.runs:
                        run.font.size = Pt(10)

        return len(prs.slides)

    def _get_placeholder(self, slide, idx: int):
        for placeholder in slide.placeholders:
            if placeholder.placeholder_format.idx == idx:
                return placeholder

        raise ValueError(f"Placeholder with idx={idx} not found")

    def _set_background_color(self, slide, color: str) -> None:
        hex_color = color.lstrip("#")
        r, g, b = tuple(int(hex_color[i:i + 2], 16) for i in (0, 2, 4))

        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(r, g, b)

    def add_uml_slide(self, prs_id: str, title: str, image_path: str) -> int:
        """
        Добавляет слайд с заголовком и PlantUML-диаграммой как изображением.
        Использует layout с картинкой по центру (аналогично add_image_content_slide).
        """
        presentation = self.store.get(prs_id)
        prs = presentation.prs

        # Используем master с layout для изображения (idx=2, layout=4 — как в image_content)
        master = prs.slide_masters[2]
        slide = prs.slides.add_slide(master.slide_layouts[4])

        title_ph = self._get_placeholder(slide, 0)
        picture_ph = self._get_placeholder(slide, 13)  # placeholder для картинки

        title_ph.text = title

        if Path(image_path).exists():
            picture_ph.insert_picture(image_path)
        else:
            raise ValueError(f"PlantUML image not found: {image_path}")

        return len(prs.slides)